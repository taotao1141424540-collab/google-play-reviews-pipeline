#!/usr/bin/env python3
"""
Build bilingual (ZH + EN) EDA conclusion deck: PPTX + PDF.

Outputs:
  reports/EDA_Conclusion_Bilingual.pptx
  reports/EDA_Conclusion_Bilingual.pdf
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT_PPTX = ROOT / "reports" / "EDA_Conclusion_Bilingual.pptx"
OUT_PDF = ROOT / "reports" / "EDA_Conclusion_Bilingual.pdf"

# Key figures (update if you re-run pipeline)
RAW_ROWS = 15_390
CLEAN_ALL_ROWS = 13_879
CLEAN_EN_ROWS = 10_072
EN_SHARE_ON_CLEAN_ALL = 0.7257


def _add_slide_title_body(prs: Presentation, title_zh: str, title_en: str, lines: list[tuple[str, str]]) -> None:
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{title_en}\n{title_zh}"
    p.font.size = Pt(22)
    p.font.bold = True

    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(9), Inches(5.2))
    btf = body.text_frame
    btf.word_wrap = True
    for i, (en, zh) in enumerate(lines):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = f"• {en}"
        para.font.size = Pt(14)
        para.space_after = Pt(4)
        para2 = btf.add_paragraph()
        para2.text = f"  {zh}"
        para2.font.size = Pt(12)
        para2.level = 0


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.5))
    t = tb.text_frame
    t.paragraphs[0].text = "Google Play Reviews — EDA Conclusion (Bilingual)"
    t.paragraphs[0].font.size = Pt(32)
    t.paragraphs[0].font.bold = True
    t.add_paragraph().text = "Google Play 评论数据 · 探索性数据分析结论（中英双语）"
    t.paragraphs[1].font.size = Pt(20)
    t.add_paragraph().text = "Sciencia AI / Mentor update · Apr 2026"
    t.paragraphs[2].font.size = Pt(14)

    _add_slide_title_body(
        prs,
        "执行摘要 Executive Summary",
        "Executive Summary · 执行摘要",
        [
            (
                f"Built an English analysis set of **{CLEAN_EN_ROWS:,}** reviews (≥10k target met).",
                f"构建英文分析集 **{CLEAN_EN_ROWS:,}** 条，满足 mentor 提出的 **10k+** 规模要求。",
            ),
            (
                "Completed EDA Sections A–E: distributions, cross-app patterns, language mix, light text, risk flags.",
                "已完成 EDA A–E：分布、跨 App 模式、语言构成、轻量词频、风险标记。",
            ),
            (
                "Heuristic flags (spam/time) are exploratory — recommend spot-checks before strong claims.",
                "刷评/时间爆发等为启发式指标，强结论前建议抽样复核。",
            ),
        ],
    )

    _add_slide_title_body(
        prs,
        "数据规模 Data volume",
        "Data volume · 数据规模",
        [
            (f"Raw rows: **{RAW_ROWS:,}** (`google_play_reviews_raw.xlsx`)", f"原始数据：**{RAW_ROWS:,}** 行"),
            (
                f"After P0 (all languages): **{CLEAN_ALL_ROWS:,}** (`clean_all_languages.xlsx`)",
                f"P0 后全语言：**{CLEAN_ALL_ROWS:,}** 行",
            ),
            (
                f"After P0 + English filter: **{CLEAN_EN_ROWS:,}** (`clean_en_only.xlsx`)",
                f"P0 + 英文子集：**{CLEAN_EN_ROWS:,}** 行",
            ),
            (
                f"English share on P0-all-lang set: ~**{EN_SHARE_ON_CLEAN_ALL:.1%}** (`is_en` / langdetect)",
                f"P0 全语言集中英文占比约 **{EN_SHARE_ON_CLEAN_ALL:.1%}**（`is_en` / 语言检测）",
            ),
        ],
    )

    _add_slide_title_body(
        prs,
        "核心分布 Core distributions (Section A)",
        "Core distributions · 核心分布（A）",
        [
            ("Rating distribution: see `eda_section_a/A1_*.png` — check skew toward 4–5★.", "评分分布：见 `eda_section_a/A1_*.png`，关注是否偏向 4–5 星。"),
            ("Review length: see `A3_*.png` — report mean / median / P90 in the written memo.", "评论长度：见 `A3_*.png`，正文报告补充均值/中位数/P90。"),
            ("Length vs rating (boxplot): `A4_*.png` — contrasts short praise vs long complaints.", "长度×评分箱线图：`A4_*.png`，可看短好评与长差评差异。"),
        ],
    )

    _add_slide_title_body(
        prs,
        "模式 Patterns (Section B)",
        "Patterns · 模式（B）",
        [
            ("Per-app means & heatmaps show category-specific rating behavior (`eda_section_b/`).", "按 App 的均值与热力图体现品类差异（`eda_section_b/`）。"),
            ("Daily volume trend flags campaign/version spikes — interpret with product calendar.", "日度评论量可对照版本/活动日历解读峰值。"),
            ("Inconsistent rating + duplicate `text_hash` tables support quality triage.", "评分-文本粗一致与重复文本表用于质量排查。"),
        ],
    )

    _add_slide_title_body(
        prs,
        "语言 Language (Section C)",
        "Language · 语言（C）",
        [
            ("Language mix chart (`eda_section_c/C1_*.png`) justifies the English subset scope.", "语言构成图（`eda_section_c/C1_*.png`）用于说明英文子集覆盖范围。"),
            ("Non-English remainder is expected — document limitations for NLP generalization.", "非英文部分客观存在，需在报告中写明对泛化的限制。"),
        ],
    )

    _add_slide_title_body(
        prs,
        "文本与风险 Text & risks (D / E)",
        "Text & risks · 文本与风险（D/E）",
        [
            ("Top tokens by overall / by star band (`eda_section_d/`) — exploratory, not topic models.", "词频与分星级词对比（`eda_section_d/`）为探索性，不等价主题模型。"),
            ("`is_spam_bot_suspect` & `is_time_anomaly` (`eda_section_e/`) — heuristic; sample before claiming abuse.", "`is_spam_bot_suspect` 与时间异常（`eda_section_e/`）为启发式，定性前需抽样。"),
        ],
    )

    _add_slide_title_body(
        prs,
        "建议 Next steps",
        "Recommendations · 建议",
        [
            ("Package: raw + cleaned exports + EDA folders + curated `reports/` spreadsheets + this deck.", "交付：raw/清洗导出、EDA 各目录、`reports/` 下整理的表格与本幻灯片。"),
            ("Optional: stricter `text_hash` dedup, refine English rule, or add topic/sentiment models if mentor agrees.", "可选：同文去重、收紧英文规则；若 mentor 同意再上主题/情感模型。"),
            ("Schedule a short review to align on which downstream task (dashboard / model) is priority.", "与 mentor 短会确认下游优先级（看板 / 模型等）。"),
        ],
    )

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    t = tb.text_frame
    t.paragraphs[0].text = "Thank you / 谢谢"
    t.paragraphs[0].font.size = Pt(36)
    t.paragraphs[0].alignment = PP_ALIGN.CENTER
    t.add_paragraph().text = "Questions & next alignment welcome.\n欢迎提问与对齐下一步。"
    t.paragraphs[1].font.size = Pt(18)
    t.paragraphs[1].alignment = PP_ALIGN.CENTER

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)


def build_pdf() -> None:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, PageBreak

    font_registered = False
    for path in (
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
    ):
        p = Path(path)
        if p.exists():
            pdfmetrics.registerFont(TTFont("UnicodeFont", str(p)))
            font_registered = True
            break

    styles = getSampleStyleSheet()
    body_en = ParagraphStyle(
        "body_en",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=11,
        leading=14,
        spaceAfter=6,
    )
    body_zh = ParagraphStyle(
        "body_zh",
        parent=styles["Normal"],
        fontName="UnicodeFont" if font_registered else "Helvetica",
        fontSize=10,
        leading=13,
        spaceAfter=10,
        leftIndent=12,
    )
    h_en = ParagraphStyle("h_en", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=16, spaceAfter=8)
    h_zh = ParagraphStyle(
        "h_zh",
        parent=styles["Heading2"],
        fontName="UnicodeFont" if font_registered else "Helvetica-Bold",
        fontSize=14,
        spaceAfter=6,
    )

    story = []
    story.append(Paragraph("Google Play Reviews — EDA Conclusion (Bilingual)", h_en))
    story.append(Paragraph("Google Play 评论数据 · 探索性数据分析结论（中英双语）", h_zh))
    story.append(Spacer(1, 0.2 * inch))

    sections = [
        (
            "1. Executive summary",
            "1. 执行摘要",
            [
                (
                    f"We deliver <b>{CLEAN_EN_ROWS:,}</b> English reviews after P0 cleaning (≥10k requirement satisfied).",
                    f"P0 清洗后英文子集 <b>{CLEAN_EN_ROWS:,}</b> 条，满足 <b>10k+</b> 要求。",
                ),
                (
                    "EDA Sections A–E cover distributions, patterns, language, token contrasts, and heuristic risk flags.",
                    "EDA A–E 覆盖分布、模式、语言、词频与启发式风险标记。",
                ),
                (
                    "Risk flags are not ground truth — manual sampling is recommended.",
                    "风险标记非金标准，建议人工抽样后再下结论。",
                ),
            ],
        ),
        (
            "2. Data volume chain",
            "2. 数据规模链路",
            [
                (f"Raw: <b>{RAW_ROWS:,}</b> rows.", f"原始：<b>{RAW_ROWS:,}</b> 行。"),
                (f"P0 all languages: <b>{CLEAN_ALL_ROWS:,}</b> rows.", f"P0 全语言：<b>{CLEAN_ALL_ROWS:,}</b> 行。"),
                (f"P0 English subset: <b>{CLEAN_EN_ROWS:,}</b> rows.", f"P0 英文：<b>{CLEAN_EN_ROWS:,}</b> 行。"),
                (
                    f"English share on P0-all-lang: ~<b>{EN_SHARE_ON_CLEAN_ALL:.1%}</b>.",
                    f"P0 全语言集中英文占比约 <b>{EN_SHARE_ON_CLEAN_ALL:.1%}</b>。",
                ),
            ],
        ),
        (
            "3. Core findings (see figures under reports/eda_section_*)",
            "3. 核心发现（图见 reports/eda_section_*）",
            [
                ("<b>A</b>: rating & length distributions; length vs rating.", "<b>A</b>：评分与长度分布；长度×评分。"),
                ("<b>B</b>: per-app differences; daily volume; inconsistency & duplicate text probes.", "<b>B</b>：跨 App 差异、日趋势、评分-文本与重复文本探索。"),
                ("<b>C</b>: language composition supports English-only modeling scope.", "<b>C</b>：语言构成支撑英文建模范围说明。"),
                ("<b>D/E</b>: token highlights; heuristic spam/time flags — use cautiously.", "<b>D/E</b>：词频亮点；刷评/时间启发式——谨慎使用。"),
            ],
        ),
        (
            "4. Deliverables & next steps",
            "4. 交付与下一步",
            [
                ("Ship: cleaned exports, EDA artifacts, curated `reports/` tables, this PDF/PPT.", "交付：清洗导出、EDA 产出、`reports/` 整理表、本 PDF/PPT。"),
                ("Next: mentor chooses downstream focus (dashboard, sentiment/topic model, more data).", "下一步：由 mentor 确定下游重点（看板、情感/主题模型、扩数据等）。"),
            ],
        ),
    ]

    for en_title, zh_title, pairs in sections:
        story.append(Paragraph(en_title, h_en))
        story.append(Paragraph(zh_title, h_zh))
        for en, zh in pairs:
            story.append(Paragraph(en, body_en))
            story.append(Paragraph(zh, body_zh))
        story.append(Spacer(1, 0.15 * inch))
        story.append(PageBreak())

    OUT_PDF.parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(
        str(OUT_PDF),
        pagesize=LETTER,
        rightMargin=inch * 0.75,
        leftMargin=inch * 0.75,
        topMargin=inch * 0.75,
        bottomMargin=inch * 0.75,
    )
    doc.build(story)


def main() -> None:
    build_pptx()
    build_pdf()
    print(f"Saved:\n- {OUT_PPTX}\n- {OUT_PDF}")


if __name__ == "__main__":
    main()
